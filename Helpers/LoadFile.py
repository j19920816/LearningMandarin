# -*- coding: UTF-8 -*-
import os
import pandas
from pptx import Presentation

import Config

encodeType = "utf8"
character = 2
radical = 3
runLine = 0


def load_dictionary_file(file_index):
    dict_dir_path = os.path.join(os.getcwd(), Config.DICT_FILE_NAME)
    data_frame = pandas.read_excel(dict_dir_path + "\\" + os.listdir(dict_dir_path)[file_index], encoding=encodeType)
    return data_frame


def load_powerpoint_file(file_name):
    presentation = Presentation(file_name)
    i = 1
    content_list = []
    while i < len(presentation.slides):
        slide = presentation.slides[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                paragraphs = shape.text_frame.paragraphs
                for paragraph in paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        if is_contains_chinese(text):
                            content_list.append(text)
        i = i + 1
    return content_list


def get_file_list(dir_name):
    return os.listdir(get_dir_path(dir_name))


def get_dir_path(dir_name):
    return os.path.join(os.getcwd(), dir_name)


def is_contains_chinese(strings):
    for char in strings:
        if '\u4e00' <= char <= '\u9fa5':
            return True
    return False
